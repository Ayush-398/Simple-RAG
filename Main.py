import tkinter as tk 
import os
import re
import ollama
from tkinter import filedialog 
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from tiktoken import encoding_for_model


EMBEDDING_MODEL = 'hf.co/CompendiumLabs/bge-base-en-v1.5-gguf'
LANGUAGE_MODEL = 'hf.co/bartowski/Llama-3.2-1B-Instruct-GGUF'

VECTOR_DB = []
text = []

root = tk.Tk()
root.withdraw()

def open_doc():
    file_path = filedialog.askopenfilename(
        title="Select a file",
        filetypes=[
            ("PDF files", "*.pdf"),
            ("PowerPoint files", "*.pptx"),
            ("Word documents", "*.docx"),
            ("All files", "*.*")
        ])
    if not file_path:
        print("No file selected.")
        return None
    if file_path.endswith(".pdf"):
        def read_pdf(file_path):
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            text = '\n'.join(text)
            text = re.sub(r'\s+', ' ', text)
            text = text.strip().lower()
            return text
        return read_pdf(file_path)
    elif file_path.endswith(".pptx"):
        def read_pptx(file_path):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            text = '\n'.join(text)
            text = re.sub(r'\s+', ' ', text)
            text = text.strip().lower()
            return text
        return read_pptx(file_path)
    elif file_path.endswith(".docx"):
         def read_docx(path):
             doc = Document(path)
             text = []
             for para in doc.paragraphs:
                 text.append(para.text)
             text = '\n'.join(text)
             text = re.sub(r'\s+', ' ', text)  
             text = text.strip().lower()
             return text
    else:
        print("Unsupported file type.")
        return None

def chunk_and_add_to_database(text):
    encoder = encoding_for_model("text-embedding-3-small")
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", " ", ""],
        length_function=lambda txt: len(encoder.encode(txt))
    )
    chunks = splitter.create_documents([text])
    for chunk in chunks:
        chunk_content = chunk.page_content
        embedding = ollama.embed(model=EMBEDDING_MODEL, input=chunk_content)['embeddings'][0]
        VECTOR_DB.append((chunk_content, embedding))
    return [chunk.page_content for chunk in chunks]
    
doc_text = open_doc()
if doc_text is None:
    print("Unsupported file type loaded")
    exit()

dataset = chunk_and_add_to_database(doc_text)
print(f"Number of chunks: {len(dataset)}")
print(f"VECTOR_DB size: {len(VECTOR_DB)}")

 
def cosine_similarity(a, b):
    dot_product = sum([x*y for x, y in zip(a, b)])
    norm_a = sum([x ** 2 for x in a]) ** 0.5
    norm_b = sum([x ** 2 for x in b]) ** 0.5
    return dot_product / (norm_a * norm_b)

def retrieve(query, top_n=3):
    query_embedding = ollama.embed(model=EMBEDDING_MODEL, input=query)['embeddings'][0]
    similarities = []
    for chunk, embedding in VECTOR_DB:
        similarity = cosine_similarity(query_embedding, embedding)
        similarities.append((chunk, similarity))
    
    similarities.sort(key=lambda x: x[1], reverse=True)
    return similarities[:top_n]

def is_context_relevant(retrieved_knowledge, query, threshold=0.4):
    return any(similarity > threshold for _, similarity in retrieved_knowledge)

input_query = input('Hi! How can I help you?: ')
retrieved_knowledge = retrieve(input_query)
print(f"Retrieved knowledge: {retrieved_knowledge}")


if not is_context_relevant(retrieved_knowledge, input_query):
    print("Sorry, I don't have enough information to answer.")
else:
    print('Answer: ')
    for chunk, similarity in retrieved_knowledge:
        print(f' - (similarity: {similarity:.2f}) {chunk}')

    prompt = f'''You are a helpful chatbot.
Use only the following pieces of context to answer the question. Don't make up any new information."
Context:
{chr(10).join([f' - {chunk}' for chunk, similarity in retrieved_knowledge])}
'''

    stream = ollama.chat(
        model=LANGUAGE_MODEL,
        messages=[
            {'role': 'system', 'content': prompt},
            {'role': 'user', 'content': input_query},
        ],
        stream=True,
    )

    print('Chatbot response:')
    for chunk in stream:
        print(chunk['message']['content'], end='', flush=True)